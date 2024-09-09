from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
import matplotlib.pyplot as plt
from io import BytesIO

# Create a new presentation
prs = Presentation()

# Function to convert LaTeX to an image and return the image path
def latex_to_image(latex_code):
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.text(0.5, 0.5, f"${latex_code}$", fontsize=12, ha='center', va='center')
    ax.axis('off')

    # Save the image to a BytesIO object
    buf = BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True, dpi=300)
    plt.close(fig)
    buf.seek(0)
    
    return buf

# Function to apply formatting to parts of a paragraph
def apply_text_formatting(paragraph, text):
    bold = False
    italic = False
    
    tokens = re.split(r'(\*\*|\*|__|_)', text)  # Split by bold (**), italics (*), and underline (__)

    for token in tokens:
        run = paragraph.add_run()  # Create a new run for each segment of text
        if token == '**':
            bold = not bold
            run.bold = bold
        elif token == '*':
            italic = not italic
            run.italic = italic
        elif token == '__':
            run.underline = True
        elif token == '_':
            run.underline = False
        else:
            run.text = token  # Assign the text to the run


# Function to add a slide with general styling
def add_slide(prs, title, content, title_font_size=32, content_font_size=24, background_color=None, title_color=RGBColor(0, 51, 102), content_color=RGBColor(0, 0, 0)):
    slide_layout = prs.slide_layouts[1]  # Slide with title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # Style for the title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(title_font_size)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = title_color

    # Change the slide background if specified
    if background_color:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background_color

    # Style for the content
    tf = content_placeholder.text_frame
    tf.clear()  # Clear any existing content

    top_position = Inches(1.5)  # Initial top position for content

    for paragraph_text in content:
        p = tf.add_paragraph()
        p.font.size = Pt(content_font_size)
        p.font.color.rgb = content_color
        
        if paragraph_text.startswith("### "):
            p.level = 0  # Subheader level (indentation)
            p.text = paragraph_text[4:].strip()  # Remove "### " from the text
            p.font.bold = True
            p.font.size = Pt(28)
        elif paragraph_text.startswith("-"):
            p.level = 1  # Bullet point
            p.text = paragraph_text[1:].strip()
        elif re.match(r'^\d+\.', paragraph_text):
            p.level = 1  # Numbered list with indent
            p.text = paragraph_text.strip()
        elif paragraph_text.startswith("$$") and paragraph_text.endswith("$$"):
            # This is for an inline equation, LaTeX style
            latex_code = paragraph_text[2:-2].strip()
            image_stream = latex_to_image(latex_code)
            
            # Insert the image in the slide
            image = slide.shapes.add_picture(image_stream, Inches(1), top_position, width=Inches(5))  # Adjust width and position as needed
            
            top_position += image.height + Inches(0.2)  # Update position for next content
        else:
            p.level = 0  # Regular paragraph
            apply_text_formatting(p, paragraph_text)  # Apply mixed formatting

# Read the text file
with open('input.md', 'r', encoding='utf-8') as file:
    lines = file.readlines()

# Process the text file
title = None
subtitle = None
author = None
current_slide_title = None
current_slide_content = []

for line in lines:
    line = line.strip()
    
    if line.startswith("# Título:"):
        title = line.replace("# Título:", "").strip()
    elif line.startswith("# Subtítulo:"):
        subtitle = line.replace("# Subtítulo:", "").strip()
    elif line.startswith("# Autor:"):
        author = line.replace("# Autor:", "").strip()
    elif line.startswith("## "):  # New slide
        if current_slide_title:
            if "Tarea" in current_slide_title:
                # Differentiate visually the Task and Task Resolution sections
                add_slide(prs, current_slide_title, current_slide_content, background_color=RGBColor(255, 230, 204), title_color=RGBColor(102, 51, 0))
            elif "Clase" in current_slide_title:
                # Start of class slide with a different style
                add_slide(prs, current_slide_title, current_slide_content, title_font_size=44, title_color=RGBColor(0, 102, 204), content_font_size=28, content_color=RGBColor(0, 0, 128))
            else:
                add_slide(prs, current_slide_title, current_slide_content)
        current_slide_title = line.replace("## ", "").strip()
        current_slide_content = []
    elif line:  # Part of the slide content
        current_slide_content.append(line)

# Add the last slide
if current_slide_title:
    if "Tarea" in current_slide_title:
        add_slide(prs, current_slide_title, current_slide_content, background_color=RGBColor(255, 230, 204), title_color=RGBColor(102, 51, 0))
    elif "Clase" in current_slide_title:
        add_slide(prs, current_slide_title, current_slide_content, title_font_size=44, title_color=RGBColor(0, 102, 204), content_font_size=28, content_color=RGBColor(0, 0, 128))
    else:
        add_slide(prs, current_slide_title, current_slide_content)

# Create the title slide
slide_layout = prs.slide_layouts[0]  # Title slide
slide = prs.slides.add_slide(slide_layout)
title_placeholder = slide.shapes.title
subtitle_placeholder = slide.placeholders[1]

title_placeholder.text = title
title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
title_placeholder.text_frame.paragraphs[0].font.bold = True
title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

subtitle_placeholder.text = f"{subtitle}\n{author}"
subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
subtitle_placeholder.text_frame.paragraphs[0].font.italic = True
subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

# Move the title slide to the beginning
xml_slides = prs.slides._sldIdLst  
slides = list(xml_slides)
xml_slides.remove(slides[-1])  # Move the title slide to the beginning
xml_slides.insert(0, slides[-1])

# Save the presentation
prs.save('output.pptx')
