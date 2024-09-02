from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
import matplotlib.pyplot as plt
from io import BytesIO
from PIL import Image

# Crear una nueva presentación
prs = Presentation()

# Function to convert LaTeX to an image and return the image path
def latex_to_image(latex_code):
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.text(0.5, 0.5, f"${latex_code}$", fontsize=12, ha='center', va='center')
    ax.axis('off')

    # Save the image to a BytesIO object
    buf = BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True, dpi=300)
    plt.close(fig)
    buf.seek(0)
    
    return buf

# Function to parse and add formatted text to a paragraph
def add_formatted_text(paragraph, text):
    parts = re.split(r'(\*\*.+?\*\*|\$\$.+?\$\$)', text)  # Split by bold or LaTeX patterns

    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            # Bold text
            run = paragraph.add_run()
            run.text = part[2:-2]  # Remove the '**'
            run.bold = True
        elif part.startswith('$$') and part.endswith('$$'):
            # LaTeX equation, render as image
            latex_code = part[2:-2].strip()
            image_stream = latex_to_image(latex_code)
            
            # Insert the image in the paragraph
            slide = paragraph._element.getparent().getparent()  # Get slide from paragraph
            img = slide.add_picture(image_stream, Inches(1), paragraph._element.top, width=Inches(5))
        else:
            # Regular text
            run = paragraph.add_run()
            run.text = part

# Function to add a slide with general styling
def add_slide(prs, title, content, title_font_size=32, content_font_size=24, background_color=None, title_color=RGBColor(0, 51, 102), content_color=RGBColor(0, 0, 0)):
    slide_layout = prs.slide_layouts[1]  # Diapositiva con título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # Estilo para el título
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(title_font_size)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = title_color

    # Cambiar el fondo de la diapositiva si se especifica
    if background_color:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background_color

    # Estilo para el contenido
    tf = content_placeholder.text_frame
    tf.clear()  # Limpiar cualquier contenido existente

    for paragraph_text in content:
        p = tf.add_paragraph()
        p.font.size = Pt(content_font_size)
        p.font.color.rgb = content_color

        if paragraph_text.startswith("### "):
            p.level = 0  # Subheader level (indentation)
            p.text = paragraph_text[4:].strip()  # Remove "### " from the text
            p.font.bold = True
            p.font.size = Pt(28)
        else:
            add_formatted_text(p, paragraph_text)  # Add formatted text

# Leer el archivo de texto
with open('contenido.txt', 'r', encoding='utf-8') as file:
    lines = file.readlines()

# Procesar el archivo de texto
title = None
subtitle = None
current_slide_title = None
current_slide_content = []

for line in lines:
    line = line.strip()
    
    if line.startswith("# Título:"):
        title = line.replace("# Título:", "").strip()
    elif line.startswith("# Subtítulo:"):
        subtitle = line.replace("# Subtítulo:", "").strip()
    elif line.startswith("# Autor:"):
        author = line.replace("# Autor:", "").strip()
    elif line.startswith("## "):  # Nueva diapositiva
        if current_slide_title:
            if "Tarea" in current_slide_title:
                # Diferenciar visualmente las secciones de Tarea y Resolución de Tarea
                add_slide(prs, current_slide_title, current_slide_content, background_color=RGBColor(255, 230, 204), title_color=RGBColor(102, 51, 0))
            elif "Clase" in current_slide_title:
                # Diapositiva de inicio de clase con un estilo distinto
                add_slide(prs, current_slide_title, current_slide_content, title_font_size=44, title_color=RGBColor(0, 102, 204), content_font_size=28, content_color=RGBColor(0, 0, 128))
            else:
                add_slide(prs, current_slide_title, current_slide_content)
        current_slide_title = line.replace("## ", "").strip()
        current_slide_content = []
    elif line:  # Parte del contenido de la diapositiva
        current_slide_content.append(line)

# Añadir la última diapositiva
if current_slide_title:
    if "Tarea" in current_slide_title:
        add_slide(prs, current_slide_title, current_slide_content, background_color=RGBColor(255, 230, 204), title_color=RGBColor(102, 51, 0))
    elif "Clase" in current_slide_title:
        add_slide(prs, current_slide_title, current_slide_content, title_font_size=44, title_color=RGBColor(0, 102, 204), content_font_size=28, content_color=RGBColor(0, 0, 128))
    else:
        add_slide(prs, current_slide_title, current_slide_content)

# Crear la diapositiva de título
slide_layout = prs.slide_layouts[0]  # Diapositiva de título
slide = prs.slides.add_slide(slide_layout)
title_placeholder = slide.shapes.title
subtitle_placeholder = slide.placeholders[1]

title_placeholder.text = title
title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
title_placeholder.text_frame.paragraphs[0].font.bold = True
title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

subtitle_placeholder.text = f"{subtitle}\n{author}"
subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
subtitle_placeholder.text_frame.paragraphs[0].font.italic = True
subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

# Mover la diapositiva de título al inicio
xml_slides = prs.slides._sldIdLst  
slides = list(xml_slides)
xml_slides.remove(slides[-1])  # Mover la diapositiva de título al inicio
xml_slides.insert(0, slides[-1])

# Guardar la presentación
prs.save('presentacion_generada.pptx')
